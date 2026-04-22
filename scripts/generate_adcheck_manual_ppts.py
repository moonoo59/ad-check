#!/usr/bin/env python3
from __future__ import annotations

import asyncio
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DIST = ROOT / "dist"
ASSETS = DIST / "ppt_assets"
BASE_URL = "http://localhost:4002"
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

W, H = Inches(13.333), Inches(7.5)

COLORS = {
    "ink": RGBColor(32, 38, 46),
    "muted": RGBColor(101, 112, 128),
    "faint": RGBColor(236, 239, 244),
    "line": RGBColor(214, 220, 229),
    "blue": RGBColor(38, 102, 214),
    "navy": RGBColor(24, 42, 71),
    "green": RGBColor(20, 132, 96),
    "orange": RGBColor(232, 130, 37),
    "red": RGBColor(208, 70, 74),
    "purple": RGBColor(112, 80, 190),
    "white": RGBColor(255, 255, 255),
    "paper": RGBColor(248, 250, 252),
}


def ensure_dirs() -> None:
    DIST.mkdir(exist_ok=True)
    ASSETS.mkdir(exist_ok=True)


async def capture_screenshots() -> dict[str, Path]:
    from playwright.async_api import async_playwright

    ensure_dirs()
    shots = {
        "login": ASSETS / "01_login.png",
        "request_new": ASSETS / "02_request_new_ad_team.png",
        "user_management": ASSETS / "03_user_management_channels.png",
        "manual": ASSETS / "04_manual_permissions.png",
        "register": ASSETS / "05_register_channels.png",
    }

    async with async_playwright() as p:
        browser = await p.chromium.launch(
            executable_path=CHROME,
            headless=True,
            args=["--no-sandbox", "--disable-dev-shm-usage"],
        )

        async def new_page():
            page = await browser.new_page(viewport={"width": 1440, "height": 920}, device_scale_factor=1)
            page.set_default_timeout(10000)
            return page

        page = await new_page()
        await page.goto(BASE_URL, wait_until="networkidle")
        await page.screenshot(path=shots["login"], full_page=False)
        await page.close()

        page = await new_page()
        await page.goto(f"{BASE_URL}/register", wait_until="networkidle")
        await page.screenshot(path=shots["register"], full_page=False)
        await page.close()

        page = await new_page()
        await page.goto(BASE_URL, wait_until="networkidle")
        await page.fill("#username", "ad1")
        await page.fill("#password", "adcheck2026")
        await page.click("button[type=submit]")
        await page.wait_for_url("**/requests/new", timeout=10000)
        await page.wait_for_selector("text=내 담당 채널")
        await page.screenshot(path=shots["request_new"], full_page=False)
        await page.close()

        page = await new_page()
        await page.goto(BASE_URL, wait_until="networkidle")
        await page.fill("#username", "admin")
        await page.fill("#password", "adcheck2026")
        await page.click("button[type=submit]")
        await page.wait_for_url("**/requests/new", timeout=10000)
        await page.goto(f"{BASE_URL}/admin/users", wait_until="networkidle")
        await page.wait_for_selector("text=담당 채널")
        await page.click("text=사용자 추가")
        await page.wait_for_selector("text=파일 전송")
        await page.screenshot(path=shots["user_management"], full_page=False)
        await page.goto(f"{BASE_URL}/manual", wait_until="networkidle")
        await page.wait_for_selector("text=파일 전송")
        await page.screenshot(path=shots["manual"], full_page=False)
        await page.close()

        await browser.close()

    return shots


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=COLORS["paper"]) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_text(slide, text: str, x, y, w, h, size=18, color=None, bold=False, align=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    r = p.runs[0]
    r.font.name = "Apple SD Gothic Neo"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color or COLORS["ink"]
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, Inches(0.62), Inches(0.42), Inches(9.8), Inches(0.52), 24, COLORS["navy"], True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.64), Inches(0.92), Inches(10.4), Inches(0.34), 10.5, COLORS["muted"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.28), Inches(12.1), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.fill.background()


def add_footer(slide, idx: int) -> None:
    add_text(slide, f"광고 증빙 요청 시스템 | {idx}", Inches(10.9), Inches(7.08), Inches(1.8), Inches(0.24), 8.5, COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_chip(slide, text: str, x, y, w, fill, color=COLORS["white"]):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.34))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Apple SD Gothic Neo"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = color
    return shp


def add_card(slide, x, y, w, h, title: str | None = None, fill=COLORS["white"], line=COLORS["line"]):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    if title:
        add_text(slide, title, x + Inches(0.25), y + Inches(0.18), w - Inches(0.5), Inches(0.32), 13, COLORS["navy"], True)
    return shp


def add_bullets(slide, items: Iterable[str], x, y, w, h, size=13, color=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Apple SD Gothic Neo"
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["ink"]
        p.space_after = Pt(8)
        p._p.get_or_add_pPr().set("marL", "220000")
        p._p.get_or_add_pPr().set("hanging", "110000")
    return box


def add_screenshot(slide, path: Path, x, y, w, h, caption: str | None = None):
    frame = add_card(slide, x, y, w, h, fill=RGBColor(255, 255, 255), line=RGBColor(198, 205, 216))
    if path.exists():
        slide.shapes.add_picture(str(path), x + Inches(0.08), y + Inches(0.08), width=w - Inches(0.16), height=h - Inches(0.16))
    else:
        add_text(slide, "화면 캡처 준비 중", x + Inches(0.3), y + Inches(0.4), w - Inches(0.6), h - Inches(0.8), 14, COLORS["muted"], align=PP_ALIGN.CENTER)
    if caption:
        add_chip(slide, caption, x + Inches(0.22), y + h - Inches(0.46), min(w - Inches(0.44), Inches(4.1)), COLORS["navy"])
    return frame


def arrow(slide, x1, y1, x2, y2, color=COLORS["blue"]):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2.2)
    line.line.end_arrowhead = True
    return line


def process_node(slide, text: str, x, y, w, h, fill, color=COLORS["white"]):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Apple SD Gothic Neo"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = color
    return shp


def make_quick_manual(shots: dict[str, Path]) -> Path:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide = blank_slide(prs)
    set_bg(slide, COLORS["navy"])
    add_text(slide, "광고 증빙 요청 시스템", Inches(0.7), Inches(0.72), Inches(8.2), Inches(0.65), 27, COLORS["white"], True)
    add_text(slide, "퀵 매뉴얼", Inches(0.72), Inches(1.34), Inches(4.2), Inches(0.45), 20, RGBColor(202, 222, 255), True)
    add_text(slide, "접속, 요청 등록, 담당 채널, 파일 전송 권한, 관리자 설정까지 한 번에 보는 배포용 안내서", Inches(0.75), Inches(2.05), Inches(6.1), Inches(0.7), 14, RGBColor(222, 230, 242))
    add_card(slide, Inches(7.45), Inches(0.76), Inches(4.95), Inches(5.55), "운영 접속 정보", RGBColor(245, 248, 252), RGBColor(245, 248, 252))
    add_bullets(slide, [
        "사용자 접속: http://adcheck.tech.net",
        "서비스 PC: http://localhost:4000",
        "앱 실행: 광고증빙요청시스템.app 더블클릭",
        "제어센터: [서버 시작] → 브라우저 자동 실행",
        "종료: [서버 중지 후 종료] 사용",
    ], Inches(7.8), Inches(1.45), Inches(4.3), Inches(2.5), 13)
    add_chip(slide, "2026.04.20 업데이트", Inches(7.8), Inches(5.45), Inches(2.25), COLORS["blue"])
    add_footer(slide, 1)

    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, "1. 로그인 및 기본 이동", "계정으로 로그인하면 요청 등록 화면이 기본으로 열립니다.")
    add_screenshot(slide, shots["login"], Inches(0.72), Inches(1.62), Inches(6.25), Inches(4.6), "로그인 화면")
    add_card(slide, Inches(7.25), Inches(1.62), Inches(5.1), Inches(4.6), "핵심 안내")
    add_bullets(slide, [
        "계정은 관리자 승인 또는 사용자 관리에서 생성합니다.",
        "최초 로그인 후 우측 상단에서 비밀번호를 변경합니다.",
        "상단 사용자 영역에서 역할과 담당 채널을 확인합니다.",
        "세션은 서버 재시작 시 초기화되므로 다시 로그인합니다.",
    ], Inches(7.55), Inches(2.25), Inches(4.45), Inches(2.7), 13)
    add_footer(slide, 2)

    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, "2. 회원가입과 담당 채널", "채널 담당자는 가입 신청 시 담당 채널을 선택하고, 승인 후 해당 채널 기준으로 요청합니다.")
    add_screenshot(slide, shots["register"], Inches(0.72), Inches(1.55), Inches(5.9), Inches(4.95), "회원가입 신청")
    add_card(slide, Inches(6.95), Inches(1.55), Inches(5.45), Inches(4.95), "승인 및 채널 배정")
    add_bullets(slide, [
        "신청자는 역할과 담당 채널을 선택해 가입을 요청합니다.",
        "관리자/대표 담당자는 신청을 승인하거나 반려합니다.",
        "담당 채널은 사용자 관리에서 다시 수정할 수 있습니다.",
        "채널명이 바뀐 경우 사용자 관리에서 담당 채널을 다시 저장합니다.",
    ], Inches(7.25), Inches(2.18), Inches(4.8), Inches(2.7), 13)
    add_footer(slide, 3)

    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, "3. 담당 채널 기반 요청 등록", "채널 담당자에게는 본인 담당 채널만 드롭다운에 표시됩니다.")
    add_screenshot(slide, shots["request_new"], Inches(0.63), Inches(1.5), Inches(7.05), Inches(5.25), "채널 담당자 요청 등록")
    add_card(slide, Inches(8.0), Inches(1.5), Inches(4.55), Inches(5.25), "등록 규칙")
    add_bullets(slide, [
        "상단에 내 담당 채널이 배지로 표시됩니다.",
        "담당 채널이 없으면 요청 등록이 비활성화됩니다.",
        "API 직접 호출도 담당 외 채널은 서버에서 차단합니다.",
        "관리자/대표 담당자는 전체 활성 채널을 선택할 수 있습니다.",
    ], Inches(8.3), Inches(2.15), Inches(3.9), Inches(2.8), 13)
    add_footer(slide, 4)

    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, "4. 파일 전송 권한", "역할과 별개로 파일 전송 권한이 있으면 탐색, 선택, 승인, 복사까지 처리합니다.")
    process_node(slide, "요청 등록", Inches(0.75), Inches(2.0), Inches(1.65), Inches(0.8), COLORS["blue"])
    process_node(slide, "자동 탐색", Inches(3.0), Inches(2.0), Inches(1.65), Inches(0.8), COLORS["green"])
    process_node(slide, "파일 선택", Inches(5.25), Inches(2.0), Inches(1.65), Inches(0.8), COLORS["orange"])
    process_node(slide, "승인/복사", Inches(7.5), Inches(2.0), Inches(1.65), Inches(0.8), COLORS["purple"])
    process_node(slide, "다운로드", Inches(9.75), Inches(2.0), Inches(1.65), Inches(0.8), COLORS["navy"])
    for x in [2.42, 4.67, 6.92, 9.17]:
        arrow(slide, x, Inches(2.4), x + Inches(0.45), Inches(2.4))
    add_card(slide, Inches(0.82), Inches(3.75), Inches(5.5), Inches(2.0), "누가 전송할 수 있나요?")
    add_bullets(slide, [
        "관리자: 항상 가능",
        "대표 담당자: 기본적으로 가능",
        "채널 담당자: 관리자에게 파일 전송 권한을 받으면 가능",
    ], Inches(1.15), Inches(4.28), Inches(4.85), Inches(1.1), 12.5)
    add_card(slide, Inches(6.75), Inches(3.75), Inches(5.5), Inches(2.0), "전송 권한 범위")
    add_bullets(slide, [
        "파일 탐색 시작/재시도",
        "후보 파일 선택 및 승인/반려",
        "복사 실행/재시도, 파일 삭제, 오전송 수정",
    ], Inches(7.08), Inches(4.28), Inches(4.85), Inches(1.1), 12.5)
    add_footer(slide, 5)

    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, "5. 관리자: 사용자/권한 관리", "사용자 목록에서 담당 채널과 파일 전송 권한을 함께 확인하고 수정합니다.")
    add_screenshot(slide, shots["user_management"], Inches(0.68), Inches(1.5), Inches(7.05), Inches(5.25), "사용자 관리")
    add_card(slide, Inches(8.05), Inches(1.5), Inches(4.45), Inches(5.25), "관리 포인트")
    add_bullets(slide, [
        "목록에 담당 채널 열이 표시됩니다.",
        "추가/수정 드로어에서 활성 채널을 체크합니다.",
        "파일 전송 체크박스로 전송 기능을 부여합니다.",
        "admin은 권한 설정과 무관하게 모든 기능을 사용할 수 있습니다.",
    ], Inches(8.35), Inches(2.15), Inches(3.85), Inches(2.8), 13)
    add_footer(slide, 6)

    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, "6. 완료 후 처리", "완료 파일은 웹에서 직접 다운로드하고, 필요 시 재전송 또는 오전송 수정을 진행합니다.")
    add_card(slide, Inches(0.78), Inches(1.58), Inches(3.7), Inches(4.55), "요청자")
    add_bullets(slide, [
        "완료 상태 확인",
        "[다운로드]로 파일 수령",
        "완료 후 1일이 지나면 자동 삭제",
        "다시 필요하면 [재전송 요청]",
    ], Inches(1.1), Inches(2.2), Inches(3.05), Inches(2.4), 13)
    add_card(slide, Inches(4.82), Inches(1.58), Inches(3.7), Inches(4.55), "전송 권한 보유자")
    add_bullets(slide, [
        "복사 실패 시 재시도",
        "완료 파일 수동 삭제",
        "오전송 수정 후 단일 항목 재탐색",
        "필요 시 반려 처리",
    ], Inches(5.15), Inches(2.2), Inches(3.05), Inches(2.4), 13)
    add_card(slide, Inches(8.85), Inches(1.58), Inches(3.7), Inches(4.55), "관리자")
    add_bullets(slide, [
        "사용자/채널 관리",
        "감사 로그 확인",
        "통계 대시보드 조회",
        "운영 앱 재생성 및 배포",
    ], Inches(9.18), Inches(2.2), Inches(3.05), Inches(2.4), 13)
    add_footer(slide, 7)

    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, "7. 자주 묻는 질문", "운영 중 많이 나오는 질문만 빠르게 정리했습니다.")
    questions = [
        ("필요한 채널이 안 보입니다", "채널 담당자는 담당 채널만 보입니다. 사용자 관리의 담당 채널 배정을 확인합니다."),
        ("파일 전송 버튼이 없습니다", "사용자 관리에서 파일 전송 권한을 부여해야 합니다."),
        ("채널 담당자도 전송할 수 있나요?", "가능합니다. 파일 전송 권한을 받으면 탐색, 선택, 승인, 복사를 수행합니다."),
        ("포트 없이 접속 가능한가요?", "서비스 PC nginx가 80번을 4000번 앱 서버로 프록시하므로 http://adcheck.tech.net 로 접속합니다."),
    ]
    y = Inches(1.6)
    for q, a in questions:
        add_card(slide, Inches(0.78), y, Inches(11.85), Inches(0.92))
        add_text(slide, q, Inches(1.05), y + Inches(0.16), Inches(3.2), Inches(0.28), 12.5, COLORS["blue"], True)
        add_text(slide, a, Inches(4.25), y + Inches(0.14), Inches(7.9), Inches(0.42), 12, COLORS["ink"])
        y += Inches(1.08)
    add_footer(slide, 8)

    out = DIST / "adcheck_deploy_quick_manual.pptx"
    prs.save(out)
    return out


def make_workflow(shots: dict[str, Path]) -> Path:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide = blank_slide(prs)
    set_bg(slide, COLORS["navy"])
    add_text(slide, "업무 흐름도", Inches(0.72), Inches(0.82), Inches(5.5), Inches(0.65), 29, COLORS["white"], True)
    add_text(slide, "기능과 권한을 연결해서 보는 광고 증빙 처리 흐름", Inches(0.75), Inches(1.55), Inches(6.7), Inches(0.5), 15, RGBColor(222, 230, 242))
    add_card(slide, Inches(0.78), Inches(2.55), Inches(11.7), Inches(2.2), "핵심 원칙", RGBColor(245, 248, 252), RGBColor(245, 248, 252))
    add_bullets(slide, [
        "채널 담당자는 담당 채널 기준으로 요청을 등록합니다.",
        "전송 기능은 역할이 아니라 파일 전송 권한으로 제어합니다.",
        "관리자는 사용자별 담당 채널과 전송 권한을 함께 관리합니다.",
    ], Inches(1.1), Inches(3.16), Inches(10.9), Inches(1.1), 16)
    add_footer(slide, 1)

    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, "1. 역할과 기능 매트릭스", "역할은 기본 범위를 정하고, 기능 권한은 세부 동작을 열어줍니다.")
    headers = ["역할", "요청 등록", "담당 채널 제한", "파일 전송", "관리자 메뉴"]
    rows = [
        ["채널 담당자", "가능", "적용", "권한 부여 시 가능", "불가"],
        ["대표 담당자", "가능", "미적용", "기본 가능", "회원 승인 일부"],
        ["관리자", "가능", "미적용", "항상 가능", "전체 가능"],
    ]
    x0, y0 = Inches(0.72), Inches(1.65)
    colw = [Inches(2.0), Inches(2.0), Inches(2.35), Inches(2.6), Inches(2.5)]
    y = y0
    for i, htxt in enumerate(headers):
        add_chip(slide, htxt, x0 + sum(colw[:i]), y, colw[i] - Inches(0.05), COLORS["navy"])
    y += Inches(0.55)
    for row in rows:
        for i, cell in enumerate(row):
            add_card(slide, x0 + sum(colw[:i]), y, colw[i] - Inches(0.05), Inches(0.7), fill=COLORS["white"])
            add_text(slide, cell, x0 + sum(colw[:i]) + Inches(0.12), y + Inches(0.2), colw[i] - Inches(0.3), Inches(0.25), 11.5, COLORS["ink"], bold=(i == 0), align=PP_ALIGN.CENTER)
        y += Inches(0.82)
    add_card(slide, Inches(0.82), Inches(5.4), Inches(11.7), Inches(0.9), "정리", RGBColor(239, 247, 255), RGBColor(190, 216, 250))
    add_text(slide, "채널 담당자는 담당 채널만 요청 등록하지만, 파일 전송 권한을 받으면 전송 담당자처럼 탐색/선택/승인/복사 작업을 수행할 수 있습니다.", Inches(1.1), Inches(5.76), Inches(11.0), Inches(0.3), 12.5, COLORS["navy"])
    add_footer(slide, 2)

    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, "2. 전체 업무 흐름", "요청 등록부터 파일 수령까지의 기본 처리 순서입니다.")
    nodes = [
        ("회원가입/로그인", COLORS["navy"]),
        ("담당 채널 확인", COLORS["blue"]),
        ("요청 등록", COLORS["blue"]),
        ("파일 탐색", COLORS["green"]),
        ("파일 선택", COLORS["orange"]),
        ("승인/복사", COLORS["purple"]),
        ("다운로드", COLORS["navy"]),
    ]
    x = Inches(0.65)
    for i, (label, color) in enumerate(nodes):
        process_node(slide, label, x, Inches(2.2), Inches(1.45), Inches(0.75), color)
        if i < len(nodes) - 1:
            arrow(slide, x + Inches(1.48), Inches(2.58), x + Inches(1.82), Inches(2.58))
        x += Inches(1.82)
    add_card(slide, Inches(0.78), Inches(3.95), Inches(3.55), Inches(1.55), "채널 담당자")
    add_bullets(slide, ["담당 채널 확인", "요청 등록", "완료 파일 다운로드"], Inches(1.05), Inches(4.42), Inches(3.0), Inches(0.8), 12)
    add_card(slide, Inches(4.75), Inches(3.95), Inches(3.55), Inches(1.55), "전송 권한 보유자")
    add_bullets(slide, ["탐색/재탐색", "파일 선택", "승인/반려/복사"], Inches(5.02), Inches(4.42), Inches(3.0), Inches(0.8), 12)
    add_card(slide, Inches(8.72), Inches(3.95), Inches(3.55), Inches(1.55), "관리자")
    add_bullets(slide, ["채널 매핑", "사용자/권한", "로그/통계"], Inches(9.0), Inches(4.42), Inches(3.0), Inches(0.8), 12)
    add_footer(slide, 3)

    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, "3. 담당 채널 요청 등록 흐름", "혼란을 줄이기 위해 화면과 서버 양쪽에서 담당 채널을 검증합니다.")
    process_node(slide, "로그인", Inches(1.0), Inches(1.85), Inches(1.7), Inches(0.75), COLORS["navy"])
    process_node(slide, "assigned_channels 조회", Inches(3.35), Inches(1.85), Inches(2.15), Inches(0.75), COLORS["blue"])
    process_node(slide, "담당 채널 드롭다운", Inches(6.15), Inches(1.85), Inches(2.15), Inches(0.75), COLORS["blue"])
    process_node(slide, "서버 검증", Inches(8.95), Inches(1.85), Inches(1.75), Inches(0.75), COLORS["orange"])
    process_node(slide, "등록 완료", Inches(11.1), Inches(1.85), Inches(1.45), Inches(0.75), COLORS["green"])
    for x in [2.75, 5.55, 8.35, 10.75]:
        arrow(slide, x, Inches(2.23), x + Inches(0.45), Inches(2.23))
    add_screenshot(slide, shots["request_new"], Inches(0.82), Inches(3.2), Inches(5.7), Inches(3.25), "담당 채널만 표시")
    add_card(slide, Inches(6.85), Inches(3.2), Inches(5.65), Inches(3.25), "예외 처리")
    add_bullets(slide, [
        "담당 채널 없음: 등록 폼 비활성화 및 관리자 문의 안내",
        "담당 외 채널 직접 전송: API에서 403 차단",
        "관리자/대표 담당자: 전체 활성 채널 선택 가능",
    ], Inches(7.15), Inches(3.88), Inches(5.0), Inches(1.55), 13)
    add_footer(slide, 4)

    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, "4. 파일 전송 권한 흐름", "파일 전송은 can_copy 권한으로 열리며, 채널 담당자도 권한이 있으면 수행할 수 있습니다.")
    process_node(slide, "요청 선택", Inches(0.9), Inches(1.78), Inches(1.7), Inches(0.72), COLORS["blue"])
    process_node(slide, "파일 탐색", Inches(3.0), Inches(1.78), Inches(1.7), Inches(0.72), COLORS["green"])
    process_node(slide, "후보 검토/선택", Inches(5.1), Inches(1.78), Inches(2.0), Inches(0.72), COLORS["orange"])
    process_node(slide, "승인 또는 반려", Inches(7.5), Inches(1.78), Inches(2.0), Inches(0.72), COLORS["purple"])
    process_node(slide, "복사/재시도", Inches(9.9), Inches(1.78), Inches(1.85), Inches(0.72), COLORS["navy"])
    for x in [2.62, 4.72, 7.12, 9.52]:
        arrow(slide, x, Inches(2.14), x + Inches(0.3), Inches(2.14))
    add_card(slide, Inches(0.82), Inches(3.15), Inches(5.65), Inches(2.55), "권한 체크")
    add_bullets(slide, [
        "admin은 항상 통과합니다.",
        "tech_team은 기본적으로 can_copy=1입니다.",
        "ad_team도 사용자 관리에서 파일 전송을 체크하면 통과합니다.",
    ], Inches(1.12), Inches(3.82), Inches(5.0), Inches(1.25), 13)
    add_card(slide, Inches(6.85), Inches(3.15), Inches(5.65), Inches(2.55), "관리자가 설정하는 곳")
    add_bullets(slide, [
        "관리자 메뉴 → 사용자 관리 → 수정",
        "담당 채널 체크",
        "파일 전송 체크 후 저장",
    ], Inches(7.15), Inches(3.82), Inches(5.0), Inches(1.25), 13)
    add_footer(slide, 5)

    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, "5. 관리자 운영 흐름", "채널과 사용자 권한을 함께 관리해야 요청 등록/전송 흐름이 자연스럽게 이어집니다.")
    add_screenshot(slide, shots["user_management"], Inches(0.72), Inches(1.45), Inches(6.6), Inches(4.9), "사용자 관리")
    add_card(slide, Inches(7.62), Inches(1.45), Inches(4.85), Inches(4.9), "체크리스트")
    add_bullets(slide, [
        "채널 매핑이 활성 상태인지 확인",
        "사용자의 담당 채널이 최신 채널명과 일치하는지 확인",
        "전송 담당이 필요한 사용자에게 파일 전송 권한 부여",
        "변경 후 사용자는 재로그인하면 최신 권한을 받음",
    ], Inches(7.92), Inches(2.12), Inches(4.25), Inches(2.5), 13)
    add_footer(slide, 6)

    out = DIST / "adcheck_workflow_overview.pptx"
    prs.save(out)
    return out


def write_markdown() -> Path:
    path = DIST / "adcheck_workflow_overview.md"
    path.write_text(
        """# 광고 증빙 요청 시스템 업무 흐름도

## 핵심 권한
- 채널 담당자(ad_team)는 담당 채널만 요청 등록할 수 있습니다.
- 파일 전송 권한(can_copy)이 있으면 채널 담당자도 파일 탐색, 파일 선택, 승인/반려, 복사 실행, 재시도, 파일 삭제, 오전송 수정을 수행할 수 있습니다.
- 관리자는 사용자 관리에서 담당 채널과 파일 전송 권한을 함께 확인/수정합니다.

## 전체 흐름
1. 사용자 로그인 및 담당 채널 확인
2. 채널 담당자가 담당 채널 기준으로 요청 등록
3. 시스템이 자동으로 파일 탐색 시작
4. 파일 전송 권한 보유자가 후보 파일 검토 및 선택
5. 승인 후 서버 로컬 전달 스토리지로 복사
6. 요청자가 웹에서 완료 파일 다운로드
7. 필요 시 재전송 요청 또는 오전송 수정

## 관리자 체크포인트
- 채널 매핑 활성 상태 확인
- 사용자별 담당 채널 최신화
- 전송 업무가 필요한 채널 담당자에게 파일 전송 권한 부여
- 권한 변경 후 재로그인 안내
""",
        encoding="utf-8",
    )
    return path


async def main() -> None:
    ensure_dirs()
    shots = await capture_screenshots()
    quick = make_quick_manual(shots)
    flow = make_workflow(shots)
    md = write_markdown()
    print(quick)
    print(flow)
    print(md)


if __name__ == "__main__":
    asyncio.run(main())
