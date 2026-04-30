#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DIST = ROOT / "dist"
STITCH = DIST / "stitch-ai-pptx-design-generator"
IMAGES = STITCH / "images"
OUT = DIST / "it_purchase_portfolio.pptx"

W, H = Inches(13.333), Inches(7.5)

COLORS = {
    "primary": RGBColor(0x00, 0x57, 0xCD),
    "primary2": RGBColor(0x0D, 0x6E, 0xFD),
    "primary_light": RGBColor(0xDA, 0xE2, 0xFF),
    "surface": RGBColor(0xF6, 0xFA, 0xFF),
    "surface2": RGBColor(0xEC, 0xF5, 0xFE),
    "card": RGBColor(0xFF, 0xFF, 0xFF),
    "line": RGBColor(0xC2, 0xC6, 0xD8),
    "line2": RGBColor(0xDB, 0xE4, 0xED),
    "ink": RGBColor(0x14, 0x1D, 0x23),
    "muted": RGBColor(0x42, 0x46, 0x55),
    "subtle": RGBColor(0x72, 0x77, 0x87),
    "green": RGBColor(0x00, 0x6D, 0x41),
    "green_light": RGBColor(0xE8, 0xF8, 0xEE),
    "orange": RGBColor(0xA5, 0x6A, 0x00),
    "orange_light": RGBColor(0xFF, 0xF1, 0xDC),
    "red": RGBColor(0xBA, 0x1A, 0x1A),
    "red_light": RGBColor(0xFF, 0xEE, 0xEC),
    "navy": RGBColor(0x00, 0x19, 0x46),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

FONT = "Apple SD Gothic Neo"
MONO_FONT = "Menlo"


def ensure_dist() -> None:
    DIST.mkdir(parents=True, exist_ok=True)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=COLORS["surface"]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_run(run, size: float, color=COLORS["ink"], bold: bool = False, font: str = FONT) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    size: float = 14,
    color=COLORS["ink"],
    bold: bool = False,
    align: PP_ALIGN | None = None,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font: str = FONT,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align or PP_ALIGN.LEFT
    p.line_spacing = 1.08
    if p.runs:
        set_run(p.runs[0], size, color, bold, font)
    return box


def add_title(slide, idx: int, title: str, subtitle: str | None = None) -> None:
    add_text(slide, f"{idx:02d}", Inches(0.62), Inches(0.42), Inches(0.6), Inches(0.28), 10, COLORS["primary"], True)
    add_text(slide, title, Inches(1.1), Inches(0.34), Inches(8.5), Inches(0.48), 22, COLORS["navy"], True)
    if subtitle:
        add_text(slide, subtitle, Inches(1.12), Inches(0.86), Inches(9.7), Inches(0.34), 10.5, COLORS["muted"])
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.28), Inches(12.1), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line2"]
    line.line.fill.background()


def add_footer(slide, idx: int) -> None:
    add_text(
        slide,
        f"IT 구매/라이선스 관리 전환 포트폴리오 | {idx:02d}",
        Inches(9.3),
        Inches(7.05),
        Inches(3.15),
        Inches(0.22),
        8.2,
        COLORS["subtle"],
        align=PP_ALIGN.RIGHT,
    )


def add_card(slide, x, y, w, h, fill=COLORS["card"], line=COLORS["line2"], radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(0.8)
    return shp


def add_chip(slide, text: str, x, y, w, fill=COLORS["primary_light"], color=COLORS["primary"], size=9.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.33))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_run(p.runs[0], size, color, True)
    return shp


def add_bullets(slide, items: Iterable[str], x, y, w, h, size: float = 12.5, color=COLORS["ink"], gap: float = 0.36):
    items = list(items)
    if not items:
        return []
    line_h = h / max(len(items), 1)
    shapes = []
    for i, item in enumerate(items):
        yy = y + Inches(i * gap)
        shapes.append(add_text(slide, f"• {item}", x, yy, w, min(line_h, Inches(0.42)), size, color))
    return shapes


def add_label_value(slide, label: str, value: str, x, y, w, accent=COLORS["primary"]):
    add_card(slide, x, y, w, Inches(0.82), COLORS["card"], COLORS["line2"])
    add_text(slide, label, x + Inches(0.18), y + Inches(0.14), w - Inches(0.36), Inches(0.2), 8.6, COLORS["subtle"], True)
    add_text(slide, value, x + Inches(0.18), y + Inches(0.38), w - Inches(0.36), Inches(0.26), 12.2, accent, True)


def add_section_tag(slide, text: str, x, y, fill=COLORS["primary"]):
    return add_chip(slide, text, x, y, Inches(1.55), fill, COLORS["white"], 9)


def add_process_node(slide, label: str, x, y, w, h, fill, text_color=COLORS["white"]):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    set_run(p.runs[0], 11, text_color, True)
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=COLORS["primary"]):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.8)
    line.line.end_arrowhead = True
    return line


def add_image(slide, path: Path, x, y, w, h, caption: str | None = None):
    add_card(slide, x, y, w, h, COLORS["white"], COLORS["line"])
    if path.exists():
        slide.shapes.add_picture(str(path), x + Inches(0.08), y + Inches(0.08), width=w - Inches(0.16), height=h - Inches(0.16))
    else:
        add_text(slide, "이미지 준비 중", x, y + h / 2 - Inches(0.15), w, Inches(0.3), 13, COLORS["subtle"], align=PP_ALIGN.CENTER)
    if caption:
        add_chip(slide, caption, x + Inches(0.22), y + h - Inches(0.48), Inches(2.2), COLORS["navy"], COLORS["white"], 8.5)


def add_three_cards(slide, cards: Sequence[tuple[str, Sequence[str], RGBColor]], y=Inches(1.68)) -> None:
    x_positions = [Inches(0.72), Inches(4.62), Inches(8.52)]
    for x, (title, bullets, accent) in zip(x_positions, cards):
        add_card(slide, x, y, Inches(3.45), Inches(4.85), COLORS["white"], COLORS["line2"])
        add_card(slide, x, y, Inches(3.45), Inches(0.12), accent, accent, radius=False)
        add_text(slide, title, x + Inches(0.28), y + Inches(0.42), Inches(2.9), Inches(0.34), 15, COLORS["navy"], True)
        add_bullets(slide, bullets, x + Inches(0.34), y + Inches(1.1), Inches(2.75), Inches(2.7), 12.2, COLORS["ink"], 0.54)


def add_statement(slide, text: str, x=Inches(0.78), y=Inches(6.15), w=Inches(11.75)) -> None:
    add_card(slide, x, y, w, Inches(0.58), COLORS["surface2"], COLORS["line"])
    add_text(slide, text, x + Inches(0.28), y + Inches(0.16), w - Inches(0.56), Inches(0.24), 12.5, COLORS["navy"], True)


def add_abstract_system_panel(slide, x, y, w, h, title: str) -> None:
    add_card(slide, x, y, w, h, COLORS["white"], COLORS["line"])
    add_text(slide, title, x + Inches(0.3), y + Inches(0.28), w - Inches(0.6), Inches(0.3), 13, COLORS["navy"], True)
    rows = [
        ("요청 접수", "검토 대기"),
        ("승인 흐름", "정책 확인"),
        ("자산 추적", "현황 동기화"),
        ("증빙 관리", "로그 보존"),
    ]
    yy = y + Inches(0.9)
    for i, (label, status) in enumerate(rows):
        add_card(slide, x + Inches(0.3), yy, w - Inches(0.6), Inches(0.58), COLORS["surface"], COLORS["line2"], radius=False)
        add_text(slide, label, x + Inches(0.48), yy + Inches(0.18), Inches(1.55), Inches(0.18), 10.8, COLORS["ink"], True)
        fill = COLORS["green_light"] if i in (0, 3) else COLORS["primary_light"]
        color = COLORS["green"] if i in (0, 3) else COLORS["primary"]
        add_chip(slide, status, x + w - Inches(1.72), yy + Inches(0.13), Inches(1.12), fill, color, 7.8)
        yy += Inches(0.72)
    add_card(slide, x + Inches(0.3), y + h - Inches(0.8), w - Inches(0.6), Inches(0.36), COLORS["navy"], COLORS["navy"], radius=False)
    add_text(slide, "익명화된 IT 구매 운영 모델", x + Inches(0.44), y + h - Inches(0.7), w - Inches(0.88), Inches(0.16), 8.5, COLORS["white"], True, align=PP_ALIGN.CENTER)


def make_deck() -> Path:
    ensure_dist()
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1. Cover
    slide = blank_slide(prs)
    set_bg(slide, COLORS["surface"])
    add_chip(slide, "PORTFOLIO", Inches(0.72), Inches(0.78), Inches(1.45), COLORS["primary"], COLORS["white"], 9)
    add_text(slide, "IT 구매/라이선스 관리 업무\n전환 포트폴리오", Inches(0.72), Inches(1.32), Inches(6.7), Inches(1.25), 29, COLORS["navy"], True)
    add_text(slide, "업무 자동화 프로젝트를 기반으로 증명하는 수요 접수, 승인 흐름, 자산 추적, 감사 대응 역량", Inches(0.76), Inches(2.78), Inches(6.2), Inches(0.76), 14.2, COLORS["muted"])
    add_label_value(slide, "포지셔닝", "업무개선 + IT 구매", Inches(0.78), Inches(4.15), Inches(2.7))
    add_label_value(slide, "문서 용도", "외부 채용/면접용", Inches(3.72), Inches(4.15), Inches(2.55))
    add_label_value(slide, "공개 정책", "강한 익명화", Inches(6.48), Inches(4.15), Inches(2.35), COLORS["green"])
    add_abstract_system_panel(slide, Inches(8.05), Inches(0.72), Inches(4.55), Inches(3.65), "IT 구매 운영 대시보드")
    add_card(slide, Inches(8.05), Inches(4.66), Inches(4.55), Inches(1.45), COLORS["white"], COLORS["line2"])
    add_text(slide, "핵심 메시지", Inches(8.34), Inches(4.94), Inches(3.9), Inches(0.28), 12.2, COLORS["primary"], True)
    add_text(slide, "구매 업무를 데이터와 프로세스로 관리하는 실무형 지원자", Inches(8.34), Inches(5.28), Inches(3.75), Inches(0.48), 15, COLORS["navy"], True)
    add_footer(slide, 1)

    # 2. Positioning
    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, 2, "전환 포지셔닝", "개발 경험을 IT 구매 업무의 운영 역량으로 번역합니다.")
    add_three_cards(
        slide,
        [
            ("개발 경험", ["현업 요청 흐름을 시스템으로 구조화", "권한·검증·로그를 코드로 구현", "운영 환경에 맞춘 배포 방식 설계"], COLORS["primary"]),
            ("구매 업무 역량", ["수요 접수와 승인 기준 정리", "라이선스·증빙 상태 추적", "부서/담당자별 책임 범위 가시화"], COLORS["green"]),
            ("면접 메시지", ["도구를 만드는 사람에서 프로세스를 관리하는 사람으로 확장", "반복 업무를 표준화하는 관점 보유", "IT 구매 실무의 데이터 기반 운영에 기여"], COLORS["orange"]),
        ],
    )
    add_statement(slide, "개발 산출물은 목적이 아니라, 구매/자산관리 업무를 정확히 이해하고 개선할 수 있음을 보여주는 근거입니다.")
    add_footer(slide, 2)

    # 3. Project one-line summary
    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, 3, "프로젝트 한 줄 요약", "수작업 요청·검수·전달 프로세스를 추적 가능한 내부 업무 시스템으로 전환했습니다.")
    add_card(slide, Inches(0.72), Inches(1.62), Inches(5.35), Inches(4.42), COLORS["white"], COLORS["line2"])
    add_text(slide, "문제 정의", Inches(1.0), Inches(1.95), Inches(4.6), Inches(0.35), 15, COLORS["red"], True)
    add_bullets(slide, ["이메일과 엑셀 중심의 수동 요청 관리", "담당자 경험에 의존한 파일 식별", "처리 상태와 승인 이력 추적 어려움", "오전송·누락 발생 시 원인 확인 지연"], Inches(1.08), Inches(2.55), Inches(4.3), Inches(2.2), 12.4, COLORS["ink"], 0.52)
    add_card(slide, Inches(7.02), Inches(1.62), Inches(5.35), Inches(4.42), COLORS["white"], COLORS["line2"])
    add_text(slide, "해결 방향", Inches(7.3), Inches(1.95), Inches(4.6), Inches(0.35), 15, COLORS["green"], True)
    add_bullets(slide, ["웹 기반 요청 접수와 상태 관리", "후보 탐색·검수·승인 흐름 표준화", "권한별 접근 범위 분리", "감사 로그로 처리 이력 보존"], Inches(7.38), Inches(2.55), Inches(4.3), Inches(2.2), 12.4, COLORS["ink"], 0.52)
    add_arrow(slide, Inches(6.25), Inches(3.85), Inches(6.86), Inches(3.85), COLORS["primary"])
    add_statement(slide, "IT 구매 업무로 치환하면, 구매 요청부터 승인·증빙·이력관리까지의 표준 운영 모델을 설계한 경험입니다.")
    add_footer(slide, 3)

    # 4. Before / After
    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, 4, "Before / After", "업무 흐름을 눈에 보이게 만들고, 담당자의 판단 기준을 시스템 안으로 옮겼습니다.")
    add_card(slide, Inches(0.72), Inches(1.62), Inches(5.7), Inches(4.7), COLORS["red_light"], RGBColor(0xFF, 0xD0, 0xCA))
    add_text(slide, "Before: 수작업", Inches(1.05), Inches(1.95), Inches(4.9), Inches(0.36), 16, COLORS["red"], True)
    add_bullets(slide, ["요청이 메일·메신저·엑셀에 분산", "담당자가 저장소를 직접 탐색", "승인 기준과 처리 이력이 개인에게 의존", "오류 발생 시 재확인과 재전달 반복"], Inches(1.08), Inches(2.65), Inches(4.8), Inches(2.3), 12.3, COLORS["ink"], 0.55)
    add_card(slide, Inches(6.9), Inches(1.62), Inches(5.7), Inches(4.7), COLORS["green_light"], RGBColor(0xB9, 0xE7, 0xC7))
    add_text(slide, "After: 시스템화", Inches(7.23), Inches(1.95), Inches(4.9), Inches(0.36), 16, COLORS["green"], True)
    add_bullets(slide, ["요청 등록과 상태 조회를 단일 화면으로 통합", "후보 탐색과 검수 기준을 구조화", "권한과 승인 절차를 시스템에서 통제", "감사 로그로 변경 이력 확인 가능"], Inches(7.26), Inches(2.65), Inches(4.8), Inches(2.3), 12.3, COLORS["ink"], 0.55)
    add_statement(slide, "핵심은 기능 개발보다 업무 기준을 정리하고, 재현 가능한 운영 프로세스로 바꾼 점입니다.")
    add_footer(slide, 4)

    # 5. IT purchase mapping
    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, 5, "IT 구매 업무와의 연결", "광고 증빙 자동화의 구성 요소를 IT 구매/라이선스 관리 언어로 재분류했습니다.")
    rows = [
        ("수요 접수", "현업 요청 등록", "구매 요청서·필요 수량 접수"),
        ("승인 흐름", "후보 검수 후 승인", "예산·정책·계약 조건 검토"),
        ("자산 추적", "전달 파일과 상태 관리", "라이선스 할당·회수·사용 현황 관리"),
        ("권한 관리", "역할별 접근 제어", "구매/검토/관리자 권한 분리"),
        ("감사 대응", "처리 이력 보존", "증빙·승인 로그·변경 이력 관리"),
    ]
    x0, y0 = Inches(0.78), Inches(1.72)
    widths = [Inches(2.0), Inches(4.55), Inches(5.65)]
    headers = ["구매 관점", "프로젝트에서 한 일", "직무 적용 의미"]
    for i, header in enumerate(headers):
        add_chip(slide, header, x0 + sum(widths[:i]), y0, widths[i] - Inches(0.08), COLORS["navy"], COLORS["white"], 9.4)
    y = y0 + Inches(0.52)
    for row in rows:
        for i, cell in enumerate(row):
            add_card(slide, x0 + sum(widths[:i]), y, widths[i] - Inches(0.08), Inches(0.68), COLORS["white"], COLORS["line2"], radius=False)
            add_text(slide, cell, x0 + sum(widths[:i]) + Inches(0.14), y + Inches(0.18), widths[i] - Inches(0.36), Inches(0.22), 10.9 if i else 11.4, COLORS["primary"] if i == 0 else COLORS["ink"], i == 0, align=PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT)
        y += Inches(0.76)
    add_statement(slide, "따라서 이 프로젝트는 IT 구매 업무의 핵심인 요청, 승인, 배정, 증빙, 감사를 직접 다룬 사례로 설명할 수 있습니다.")
    add_footer(slide, 5)

    # 6. Process design
    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, 6, "업무 프로세스 설계", "요청부터 이력 관리까지 담당자와 시스템의 책임 범위를 명확히 나눴습니다.")
    steps = [
        ("요청\n등록", COLORS["primary"]),
        ("후보\n탐색", COLORS["green"]),
        ("검수", COLORS["orange"]),
        ("승인", RGBColor(0x52, 0x5A, 0xC7)),
        ("전달", COLORS["navy"]),
        ("이력\n관리", COLORS["primary2"]),
    ]
    x = Inches(0.9)
    for i, (label, color) in enumerate(steps):
        add_process_node(slide, label, x, Inches(2.05), Inches(1.35), Inches(0.86), color)
        if i < len(steps) - 1:
            add_arrow(slide, x + Inches(1.4), Inches(2.48), x + Inches(1.78), Inches(2.48))
        x += Inches(1.9)
    add_three_cards(
        slide,
        [
            ("요청자", ["필요 항목 입력", "처리 상태 확인", "완료 결과 수령"], COLORS["primary"]),
            ("검수/전달 담당", ["후보 파일 확인", "승인·반려 판단", "전달 결과 재확인"], COLORS["orange"]),
            ("관리자", ["사용자와 권한 관리", "채널/분류 기준 관리", "로그와 통계 확인"], COLORS["green"]),
        ],
        y=Inches(3.55),
    )
    add_footer(slide, 6)

    # 7. Dashboard
    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, 7, "라이선스/자산관리 관점 대시보드", "현황을 한 화면에서 파악하고 구매 판단으로 연결하는 구조를 제안합니다.")
    add_image(slide, IMAGES / "slide7_dashboard.png", Inches(0.72), Inches(1.55), Inches(6.5), Inches(5.25), "Dashboard")
    add_card(slide, Inches(7.55), Inches(1.55), Inches(4.92), Inches(5.25), COLORS["white"], COLORS["line2"])
    add_text(slide, "구매 업무 적용 포인트", Inches(7.88), Inches(1.95), Inches(4.2), Inches(0.32), 15, COLORS["navy"], True)
    add_bullets(slide, ["부서별 라이선스 집행률 확인", "추가 구매 요청과 승인 대기 건 추적", "만료·갱신 예정 항목 사전 알림", "사용률이 낮은 자산 회수 후보 식별", "정책 위반이나 미승인 사용 리스크 확인"], Inches(7.92), Inches(2.65), Inches(4.05), Inches(2.9), 12.2, COLORS["ink"], 0.49)
    add_statement(slide, "대시보드는 예쁜 화면이 아니라, 구매 우선순위와 리스크를 빠르게 판단하기 위한 운영 장치입니다.")
    add_footer(slide, 7)

    # 8. My role
    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, 8, "나의 역할", "요구사항 파악부터 정책 설계, 구현, 운영 문서화까지 전 과정을 단독으로 수행했습니다.")
    add_card(slide, Inches(7.18), Inches(1.55), Inches(5.12), Inches(4.1), COLORS["white"], COLORS["line2"])
    add_text(slide, "역할 범위", Inches(7.52), Inches(1.88), Inches(4.35), Inches(0.32), 15, COLORS["navy"], True)
    add_bullets(
        slide,
        [
            "현업 인터뷰와 문제 구조화",
            "요청·승인·전달 정책 정의",
            "권한별 화면과 API 흐름 구현",
            "운영 매뉴얼과 개선 과제 정리",
        ],
        Inches(7.58),
        Inches(2.55),
        Inches(4.05),
        Inches(1.9),
        12.2,
        COLORS["ink"],
        0.5,
    )
    add_card(slide, Inches(7.58), Inches(4.82), Inches(4.16), Inches(0.44), COLORS["surface2"], COLORS["line"])
    add_text(slide, "직무 전환 키워드: 수요관리 · 승인정책 · 자산운영 · 감사대응", Inches(7.78), Inches(4.96), Inches(3.75), Inches(0.16), 8.7, COLORS["primary"], True)
    role_cards = [
        ("요구사항 분석", "현업 요청 방식과 반복 오류를 정리하고 표준 입력 항목으로 전환"),
        ("정책 설계", "역할, 담당 범위, 전송 권한, 승인 절차를 명확히 분리"),
        ("시스템 구현", "요청·탐색·검수·전달·다운로드·로그 흐름을 웹 시스템으로 구현"),
        ("운영 문서화", "관리자 체크포인트와 사용자 흐름을 문서와 매뉴얼로 정리"),
    ]
    y = Inches(1.58)
    for title, body in role_cards:
        add_card(slide, Inches(0.82), y, Inches(5.75), Inches(0.92), COLORS["white"], COLORS["line2"])
        add_text(slide, title, Inches(1.1), y + Inches(0.18), Inches(1.55), Inches(0.24), 12.4, COLORS["primary"], True)
        add_text(slide, body, Inches(2.75), y + Inches(0.16), Inches(3.45), Inches(0.36), 11.6, COLORS["ink"])
        y += Inches(1.03)
    add_statement(slide, "IT 구매 직무에서는 이 경험을 수요관리, 승인정책, 자산운영, 감사대응 역량으로 전환해 설명합니다.")
    add_footer(slide, 8)

    # 9. Risk management
    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, 9, "리스크 관리", "구매 업무에서 중요한 통제, 감사, 오류 방지 관점을 시스템 설계에 반영했습니다.")
    add_three_cards(
        slide,
        [
            ("권한 분리", ["역할별 화면과 API 접근 제어", "담당 범위 기반 요청 제한", "관리자 권한 변경 지점 통제"], COLORS["primary"]),
            ("감사 가능성", ["요청·승인·반려·삭제 이력 기록", "사용자와 시점 기준 추적", "운영 중 원인 확인 가능한 로그 구조"], COLORS["green"]),
            ("오류 방지", ["경로 이탈 방어", "로그인 시도 제한", "시간대와 날짜 경계 정규화"], COLORS["orange"]),
        ],
        y=Inches(1.62),
    )
    add_card(slide, Inches(0.82), Inches(5.55), Inches(11.65), Inches(0.82), COLORS["white"], COLORS["line2"])
    add_text(slide, "구매 업무 적용", Inches(1.12), Inches(5.78), Inches(1.8), Inches(0.24), 12.2, COLORS["primary"], True)
    add_text(slide, "구매 요청 권한, 승인 근거, 계약 증빙, 라이선스 할당 이력까지 추적 가능한 운영 기준으로 확장할 수 있습니다.", Inches(2.75), Inches(5.76), Inches(9.2), Inches(0.28), 12.2, COLORS["ink"])
    add_footer(slide, 9)

    # 10. Operational outcomes
    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, 10, "운영 개선 성과", "검증되지 않은 숫자 대신 실제로 확인 가능한 변화 중심으로 정리했습니다.")
    outcomes = [
        ("추적 가능성", "요청 상태와 처리 이력을 시스템에서 확인"),
        ("재작업 감소", "오류 발생 시 해당 항목만 재검수·재전달 가능"),
        ("상태 가시화", "요청자와 담당자가 같은 기준으로 진행 상황 공유"),
        ("수동 전달 리스크 완화", "담당자의 기억과 로컬 파일 관리 의존도 축소"),
    ]
    positions = [(0.78, 1.65), (6.74, 1.65), (0.78, 4.02), (6.74, 4.02)]
    for (title, body), (xv, yv) in zip(outcomes, positions):
        add_card(slide, Inches(xv), Inches(yv), Inches(5.45), Inches(1.65), COLORS["white"], COLORS["line2"])
        add_text(slide, title, Inches(xv + 0.32), Inches(yv + 0.28), Inches(4.75), Inches(0.32), 16, COLORS["primary"], True)
        add_text(slide, body, Inches(xv + 0.32), Inches(yv + 0.82), Inches(4.55), Inches(0.44), 12.6, COLORS["ink"])
    add_statement(slide, "면접에서는 과장된 효율 수치보다, 어떤 리스크가 어떤 운영 장치로 줄었는지를 중심으로 설명합니다.")
    add_footer(slide, 10)

    # 11. IT purchase scenarios
    slide = blank_slide(prs)
    set_bg(slide)
    add_title(slide, 11, "IT 구매 직무 적용 시나리오", "현재 경험을 SW 라이선스 구매와 IT 자산 운영 업무로 확장하는 방식입니다.")
    scenarios = [
        ("SW 라이선스 구매 요청", ["현업 수요 접수", "수량·사용 목적 확인", "승인 상태 추적"]),
        ("부서별 사용 현황", ["할당 현황 대시보드", "미사용 자산 회수 후보", "추가 구매 필요성 판단"]),
        ("갱신/만료 관리", ["계약 만료 예정 알림", "갱신 우선순위 정리", "비용 누수 예방"]),
        ("증빙 관리", ["승인 로그 보존", "구매 근거 연결", "감사 대응 자료화"]),
    ]
    x_positions = [Inches(0.82), Inches(3.88), Inches(6.94), Inches(10.0)]
    for i, (title, bullets) in enumerate(scenarios):
        x = x_positions[i]
        add_card(slide, x, Inches(1.78), Inches(2.48), Inches(4.3), COLORS["white"], COLORS["line2"])
        add_chip(slide, f"0{i+1}", x + Inches(0.28), Inches(2.08), Inches(0.62), COLORS["primary_light"], COLORS["primary"], 9)
        add_text(slide, title, x + Inches(0.28), Inches(2.62), Inches(1.92), Inches(0.46), 13.2, COLORS["navy"], True)
        add_bullets(slide, bullets, x + Inches(0.28), Inches(3.42), Inches(1.92), Inches(1.4), 11.6, COLORS["ink"], 0.42)
    add_statement(slide, "이직 후에는 구매 업무를 단순 처리하지 않고, 요청-승인-자산-증빙이 연결되는 운영 체계로 관리하겠습니다.")
    add_footer(slide, 11)

    # 12. Closing
    slide = blank_slide(prs)
    set_bg(slide, COLORS["navy"])
    add_chip(slide, "CLOSING", Inches(0.78), Inches(0.82), Inches(1.35), COLORS["primary2"], COLORS["white"], 9)
    add_text(slide, "구매 업무를 데이터와 프로세스로\n관리하는 실무형 지원자", Inches(0.78), Inches(1.42), Inches(7.1), Inches(1.28), 27, COLORS["white"], True)
    add_text(slide, "반복 업무를 표준화하고, 승인과 증빙의 흐름을 추적 가능하게 만들며, IT 자산 운영 리스크를 줄이는 방향으로 기여하겠습니다.", Inches(0.82), Inches(3.05), Inches(6.7), Inches(0.72), 14, RGBColor(0xE9, 0xF2, 0xFB))
    closing_cards = [
        ("업무 이해", "현업 요청과 담당자 운영 흐름을 구조화"),
        ("프로세스 설계", "권한, 승인, 검수, 감사 기준을 명확화"),
        ("실행력", "작동하는 시스템과 문서로 완성"),
    ]
    y = Inches(4.42)
    for title, body in closing_cards:
        add_card(slide, Inches(0.82), y, Inches(6.45), Inches(0.58), RGBColor(0x29, 0x31, 0x38), RGBColor(0x42, 0x46, 0x55))
        add_text(slide, title, Inches(1.08), y + Inches(0.16), Inches(1.4), Inches(0.2), 10.8, RGBColor(0xB1, 0xC5, 0xFF), True)
        add_text(slide, body, Inches(2.52), y + Inches(0.15), Inches(4.25), Inches(0.22), 10.8, COLORS["white"])
        y += Inches(0.72)
    add_abstract_system_panel(slide, Inches(8.15), Inches(1.0), Inches(4.35), Inches(3.5), "Portfolio Summary")
    add_text(slide, "외부 제출용 익명화 문서", Inches(8.3), Inches(4.85), Inches(3.9), Inches(0.24), 10.8, RGBColor(0xB1, 0xC5, 0xFF), True, align=PP_ALIGN.CENTER)
    add_text(slide, "회사명, 실제 채널명, 내부망 주소, 실제 스토리지명 미포함", Inches(8.25), Inches(5.24), Inches(3.9), Inches(0.44), 11, RGBColor(0xE9, 0xF2, 0xFB), align=PP_ALIGN.CENTER)
    add_text(slide, "IT 구매/라이선스 관리 전환 포트폴리오 | 12", Inches(9.3), Inches(7.05), Inches(3.15), Inches(0.22), 8.2, RGBColor(0xE9, 0xF2, 0xFB), align=PP_ALIGN.RIGHT)

    prs.save(OUT)
    return OUT


def main() -> None:
    out = make_deck()
    print(out)


if __name__ == "__main__":
    main()
